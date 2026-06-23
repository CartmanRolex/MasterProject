"""Build presentation.pptx from the report figures.

Editable PowerPoint deck for "Enabling Recovery in VLA-Based Robotic
Manipulation". Slide order mirrors the LaTeX report (Intro -> Methods -> Results
-> Discussion); nothing is shown before it is defined. Style: terse /
keyword-driven, numeric results live in tables, not prose.

Run with the project venv:
    MasterProject/.venv/Scripts/python.exe report/slides/build_pptx.py
"""

from __future__ import annotations

import os

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.normpath(os.path.join(HERE, "..", "figures"))
ASSET_DIR = os.path.join(HERE, "assets")
OUT_PPTX = os.path.join(HERE, "presentation.pptx")

os.makedirs(ASSET_DIR, exist_ok=True)

# ---- palette ---------------------------------------------------------------
EPFL_RED = RGBColor(0xE0, 0x40, 0x4F)
INK = RGBColor(0x20, 0x24, 0x2C)
SLATE = RGBColor(0x53, 0x5A, 0x66)
LIGHT = RGBColor(0xF4, 0xF5, 0xF7)
MIDGREY = RGBColor(0xD8, 0xDC, 0xE2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_W = SLIDE_W - 2 * MARGIN


# ---- image helpers ---------------------------------------------------------
def rasterize(name_pdf: str, dpi: int = 200) -> str:
    src = os.path.join(FIG_DIR, name_pdf)
    out = os.path.join(ASSET_DIR, os.path.splitext(name_pdf)[0] + ".png")
    doc = fitz.open(src)
    pix = doc[0].get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72), alpha=False)
    pix.save(out)
    doc.close()
    return out


def asset(name: str) -> str:
    if name.lower().endswith(".pdf"):
        return rasterize(name)
    return os.path.join(FIG_DIR, name)


def img_size(path: str):
    with Image.open(path) as im:
        return im.size


def place_image(slide, path, box_l, box_t, box_w, box_h):
    iw, ih = img_size(path)
    ar = iw / ih
    box_ar = box_w / box_h
    if ar > box_ar:
        w, h = box_w, int(box_w / ar)
    else:
        h, w = box_h, int(box_h * ar)
    l = int(box_l + (box_w - w) / 2)
    t = int(box_t + (box_h - h) / 2)
    return slide.shapes.add_picture(path, l, t, width=w, height=h)


# ---- slide scaffolding -----------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_title(slide, text, sub=None):
    bar = slide.shapes.add_shape(1, MARGIN, Inches(0.5), Inches(0.09), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = EPFL_RED
    bar.line.fill.background()
    tf = add_textbox(slide, MARGIN + Inches(0.22), Inches(0.42), CONTENT_W - Inches(0.22), Inches(1.0))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = INK
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(14)
        r2.font.italic = True
        r2.font.color.rgb = EPFL_RED


def add_bullets(slide, bullets, l, t, w, h, size=18):
    tf = add_textbox(slide, l, t, w, h)
    first = True
    for item in bullets:
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(7)
        run = p.add_run()
        run.text = ("•  " if level == 0 else "–  ") + text
        run.font.size = Pt(size if level == 0 else size - 3)
        run.font.bold = level == 0
        run.font.color.rgb = INK if level == 0 else SLATE
    return tf


def add_stat_row(slide, stats, top, height=Inches(1.45), color=EPFL_RED):
    n = len(stats)
    gap = Inches(0.25)
    cell_w = int((CONTENT_W - gap * (n - 1)) / n)
    for i, (num, label) in enumerate(stats):
        l = int(MARGIN + i * (cell_w + gap))
        box = slide.shapes.add_shape(1, l, top, cell_w, height)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = color
        box.line.width = Pt(1.25)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(13)
        r2.font.color.rgb = INK


def add_label(slide, text, l, t, w, size=15, color=INK, bold=True):
    tf = add_textbox(slide, l, t, w, Inches(0.35))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_table(slide, rows, l, t, w, h, font=12, highlight_col=None):
    """rows[0] is the header. Optionally highlight a body column (1-based)."""
    n_rows, n_cols = len(rows), len(rows[0])
    gtable = slide.shapes.add_table(n_rows, n_cols, l, t, w, h).table
    gtable.first_row = False
    gtable.horz_banding = False
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = gtable.cell(r, c)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = INK
            elif highlight_col is not None and c == highlight_col:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFD, 0xE9, 0xEB)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(font)
            run.font.bold = (r == 0) or (highlight_col is not None and c == highlight_col)
            run.font.color.rgb = WHITE if r == 0 else INK
    return gtable


# ---- slide templates -------------------------------------------------------
def text_image_slide(title, bullets, images, sub=None, img_ratio=0.46, bullet_size=18):
    s = add_slide()
    add_bg(s)
    set_title(s, title, sub)
    top = Inches(1.7)
    h = SLIDE_H - Inches(0.4) - top
    text_w = int(CONTENT_W * (1 - img_ratio)) - Inches(0.2)
    add_bullets(s, bullets, MARGIN, top, text_w, h, size=bullet_size)
    img_l = MARGIN + text_w + Inches(0.3)
    img_w = SLIDE_W - MARGIN - img_l
    cell_h = int(h / len(images))
    for i, im in enumerate(images):
        place_image(s, asset(im), img_l, int(top + i * cell_h), img_w, cell_h - Inches(0.1))
    return s


def text_only_slide(title, bullets, sub=None, bullet_size=20):
    s = add_slide()
    add_bg(s)
    set_title(s, title, sub)
    add_bullets(s, bullets, MARGIN, Inches(1.9), CONTENT_W, SLIDE_H - Inches(2.4), size=bullet_size)
    return s


def big_image_slide(title, images, caption=None, sub=None, row=True):
    s = add_slide()
    add_bg(s)
    set_title(s, title, sub)
    top = Inches(1.7)
    bottom = SLIDE_H - (Inches(0.85) if caption else Inches(0.4))
    h = bottom - top
    n = len(images)
    if row:
        cell_w = int(CONTENT_W / n)
        for i, im in enumerate(images):
            place_image(s, asset(im), int(MARGIN + i * cell_w), top, cell_w - Inches(0.15), h)
    else:
        place_image(s, asset(images[0]), MARGIN, top, CONTENT_W, h)
    if caption:
        tf = add_textbox(s, MARGIN, bottom + Inches(0.05), CONTENT_W, Inches(0.7))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = caption
        r.font.size = Pt(15)
        r.font.italic = True
        r.font.color.rgb = SLATE
    return s


PRELIM = "Preliminary results — numbers may still change"

# === 1 — Title ==============================================================
s = add_slide()
add_bg(s, INK)
try:
    place_image(s, asset("logo.png"), Inches(5.4), Inches(0.7), Inches(2.5), Inches(1.4))
except Exception:
    pass
tf = add_textbox(s, Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.2), MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Enabling Recovery in VLA-Based Robotic Manipulation"
r.font.size = Pt(40)
r.font.bold = True
r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Master's Project"
r2.font.size = Pt(20)
r2.font.color.rgb = EPFL_RED
tf3 = add_textbox(s, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.6), MSO_ANCHOR.MIDDLE)
for line, sz in [
    ("Gal Pascual", 22),
    ("EPFL  •  CREATE Lab (Prof. Josie Hughes)  •  Sycamore Lab", 16),
    ("Supervisors: A. Schlaginhaufen, C. Pan, T. Ni", 14),
]:
    pp = tf3.paragraphs[0] if line == "Gal Pascual" else tf3.add_paragraph()
    pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run()
    rr.text = line
    rr.font.size = Pt(sz)
    rr.font.color.rgb = LIGHT

# === 2 — Problem (Intro) ====================================================
text_only_slide(
    "The Problem",
    [
        "Long-horizon task = grasp → lift → place, in order.",
        "Policies trained only on success → little data on recovering from failure.",
        "No built-in check that a subtask actually worked.",
        "Stagnation: retries one orange until timeout, ignores the others.",
        "Hypothesis: a failed grasp can leave the orange in an unseen state → confuses the policy.",
        "Idea: fix the structure around the policy, not just the policy.",
    ],
)

# === 3 — Hypothesis & Contributions (Intro) =================================
text_only_slide(
    "Hypothesis & Contributions",
    [
        "Hypothesis: decompose into language subtasks + an orchestrator that checks each outcome.",
        "Subtask dataset — targeted, relative-position grasp instructions.",
        "Orchestrator — monitors outcomes; retry / reset / redirect.",
        "Probe: is the subtask policy just data-limited? → autonomous data generation.",
    ],
)

# === 4 — Setup / Environment (Methods) ======================================
text_image_slide(
    "Setup",
    [
        "SO-101 arm, Isaac Sim kitchen.",
        "Task: 3 oranges → plate.",
        "Policy: SmolVLA (Vision-Language-Action).",
        ("Sees: 6 joints + front & wrist cameras.", 1),
        ("No object coordinates.", 1),
        "Privileged state (poses, forces) → orchestrator & eval only.",
    ],
    ["environment_top_view.png", "environment_plate_closeup.png"],
    img_ratio=0.42,
)

# === 5 — Subtask Decomposition (Methods) ====================================
big_image_slide(
    "Subtask Decomposition",
    ["grasp.png", "pick.png", "place.png"],
    caption="Full task → GRASP → LIFT → PLACE.  Each subtask is a checkpoint: verify the physical outcome before the next step.  Grasp targets one orange by relative position (left / middle / right).",
    row=True,
)

# === 6 — Teleop Dataset (Methods) ===========================================
s = add_slide()
add_bg(s)
set_title(s, "Custom Teleop Dataset")
add_stat_row(s, [("846", "manually teleoperated\nsubtask episodes")], top=Inches(1.7),
             height=Inches(1.4))
# make the single stat box narrower & left-aligned with bullets to its right
add_bullets(
    s,
    [
        "Built from scratch: Xbox gamepad → SO-101 in Isaac.",
        "1 episode = 1 successful subtask, labelled with its instruction.",
        "Randomised pick order → forces grounding of the position instruction.",
        "“Freeze tail” (20 frames) → policy learns to wait at a boundary instead of drifting.",
        "Trained SmolVLA as a subtask executor (GRASP / LIFT / PLACE).",
    ],
    MARGIN, Inches(3.5), CONTENT_W, Inches(3.4), size=19,
)

# === 7 — Spatial context-dependence (motivates the reset) ===================
text_only_slide(
    "Language is Spatially Context-Dependent",
    [
        "Gripper near orange A → “grasp the right orange” is ignored.",
        "Target switching only works when the arm is far from all oranges.",
        "→ Before any target change, the arm must first move away.",
        "This is exactly why the orchestrator needs a scripted reset (next slide).",
    ],
    sub="The key behavioural insight — it motivates the orchestrator's reset",
)

# === 8 — Orchestrator (Methods) — bigger figure =============================
s = add_slide()
add_bg(s)
set_title(s, "Orchestrator: Three Mechanisms")
top = Inches(1.65)
h = SLIDE_H - Inches(0.4) - top
text_w = int(CONTENT_W * 0.40)
add_bullets(
    s,
    [
        "FSM wraps the VLA; VLA does only GRASP / LIFT / PLACE.",
        ("Local retry = recovery — slip → re-GRASP same orange.", 1),
        ("Home reset = precondition — scripted move home before a target change.", 1),
        ("Redirection = give up — grasp timeout → try a different orange.", 1),
        "Checks use privileged state (contact, poses).",
    ],
    MARGIN, top, text_w, h, size=16,
)
place_image(s, asset("orchestrator_decision_flow.pdf"),
            int(MARGIN + text_w + Inches(0.25)), top,
            int(SLIDE_W - MARGIN - (MARGIN + text_w + Inches(0.25))), h)

# === 9 — Subtask Failure Modes (Methods, after orchestrator) ================
s = add_slide()
add_bg(s)
set_title(s, "Failures vs. Timeouts")
add_label(s, "Physical failure  (the orange is actually lost)", MARGIN, Inches(1.7),
          CONTENT_W, size=18, color=EPFL_RED)
add_bullets(
    s,
    [
        ("LIFT — orange slips before reaching the plate.", 1),
        ("PLACE — orange dropped outside the plate.", 1),
        ("GRASP — no instantaneous failure event: the gripper just never closes on it.", 1),
    ],
    MARGIN, Inches(2.15), CONTENT_W, Inches(2.0), size=18,
)
add_label(s, "Timeout  (only exists with the orchestrator)", MARGIN, Inches(4.35),
          CONTENT_W, size=18, color=EPFL_RED)
add_bullets(
    s,
    [
        ("A step budget per subtask — so we never stay stuck.", 1),
        ("It is a detection trigger to react (retry / reset / redirect), not a failure itself.", 1),
        ("Monotask has no timeout → a failure is never caught mid-episode.", 1),
    ],
    MARGIN, Inches(4.8), CONTENT_W, Inches(2.2), size=18,
)

# === 10 — Autonomous data generation + dataset composition (Methods) ========
text_image_slide(
    "Autonomous Data Generation",
    [
        "Goal: test if the subtask policy is just data-limited (not the recovery goal itself).",
        "Orchestrator runs on its own; keeps only successful subtasks.",
        "Redirection salvages data from otherwise-wasted scenes.",
        ("+414 Auto episodes merged with 846 Teleop → 1260 (Teleop+Auto).", 1),
        "Auto episodes are shorter / less dispersed (policy's own states).",
    ],
    ["dataset_composition.pdf"],
    img_ratio=0.5,
    bullet_size=17,
)

# === 11 — Results: Full-Task Outcomes =======================================
text_image_slide(
    "Results: Full-Task Outcomes",
    [
        "Same 100 seeded scenes per model.",
        "Baseline bars = 60 public full-task demos (reference).",
        "Subtasking → more partial progress than end-to-end.",
        "Subtask runs rarely fail completely — almost always ≥ 1 orange.",
        "But partial progress ≠ finishing all three (yet).",
    ],
    ["orange_outcome_distribution.pdf"],
    img_ratio=0.52,
    sub=PRELIM,
    bullet_size=17,
)

# === 12 — Results: Lift & Place Failures (tables) ===========================
s = add_slide()
add_bg(s)
set_title(s, "Results: Lift & Place Failures", sub=PRELIM)
add_label(s, "After a dropped LIFT — does it recover?", MARGIN, Inches(1.7), CONTENT_W, size=17)
add_table(
    s,
    [
        ["Training set", "Formulation", "Drop recovered", "Re-engaged same orange"],
        ["Teleop", "Monotask", "26%", "50%"],
        ["Teleop", "Subtask", "73%", "84%"],
        ["Teleop+Auto", "Monotask", "14%", "25%"],
        ["Teleop+Auto", "Subtask", "39%", "63%"],
    ],
    MARGIN, Inches(2.1), CONTENT_W, Inches(2.1), font=14, highlight_col=2,
)
add_label(s, "PLACE success (orchestrated subtask models)", MARGIN, Inches(4.55), CONTENT_W, size=17)
add_table(
    s,
    [
        ["Training set", "Place success", "Drops", "Subtask timeouts"],
        ["Teleop", "83%", "33", "3"],
        ["Teleop+Auto", "76%", "27", "20"],
    ],
    MARGIN, Inches(4.95), CONTENT_W, Inches(1.4), font=14, highlight_col=1,
)
add_bullets(
    s,
    [("Lift itself is similar across formulations — decomposition makes a failed lift recoverable (local retry).", 0)],
    MARGIN, Inches(6.55), CONTENT_W, Inches(0.8), size=14,
)

# === 13 — Results: Target Obedience =========================================
text_image_slide(
    "Results: Target Obedience",
    [
        "Does it grasp the orange it was told to?",
        "~3 in 4 grasps obey; better as the scene clears.",
        "Misgrabs are systematic: closes on the orange just inward.",
        ("Most common label “right” = worst obeyed.", 1),
        ("More data ≠ better obedience.", 1),
        "→ A grounding gap, not a coverage gap.",
    ],
    ["grasp_obedience_confusion.pdf"],
    img_ratio=0.52,
    sub=PRELIM,
    bullet_size=16,
)

# === 14 — Discussion & Limitations ==========================================
text_only_slide(
    "Discussion & Limitations",
    [
        "Decomposition creates real intervention points — failures caught, not ignored.",
        "But can't fix a weak GRASP policy — the bottleneck, worse as the scene fills.",
        "Auto data did not help (full-task dropped) — bottleneck is grasp quality, not volume.",
        "Coverage gap: trained on clean isolated subtasks, never on messy full runs.",
        "Success checks use privileged state — real world needs perception.",
    ],
)

# === 15 — Conclusion & Future Work ==========================================
text_only_slide(
    "Conclusion & Future Work",
    [
        "Recovery = an execution-level problem, not just a policy choice.",
        "Orchestrator enables retry / reset / redirection — and makes recovery measurable.",
        "Orchestration alone ≠ robustness while the grasp policy is weak.",
        "Next: better GRASP / PLACE in late scenes; collect full, imperfect runs — before more auto data.",
    ],
    sub="A structured reflection, including a useful negative result",
)

prs.save(OUT_PPTX)
print(f"Wrote {OUT_PPTX}")
print(f"Slides: {len(prs.slides._sldIdLst)}")
